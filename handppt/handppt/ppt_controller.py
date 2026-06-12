"""
PPT 控制器模块
功能：
1. 通过 python-pptx 解析 PPTX 文件，建立 页码→文字列表+位置 映射表
2. 键盘模拟翻页（PageDown/PageUp，兼容 PowerPoint / WPS / 任何放映软件）
3. COM 用于获取实际页码 + 幻灯片内高亮
4. 提供当前页文字列表供语音匹配
"""

import math
import os
import re
import threading
import time
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    Presentation = None

import config

# COM 支持（可选）
try:
    import win32com.client
    import pythoncom
    import win32api
    import win32con
    _HAS_COM = True
except ImportError:
    _HAS_COM = False
    win32api = None
    win32con = None


# 虚拟键码常量（避免 COM 不可用时引用失败）
VK_NEXT = 0x22      # PageDown
VK_PRIOR = 0x21     # PageUp
VK_HOME = 0x24
VK_END = 0x23
VK_ESCAPE = 0x1B
VK_B = 0x42


def _press_key(vk_code: int):
    """发送单个按键"""
    if win32api is None:
        return
    try:
        extra = 0
        win32api.keybd_event(vk_code, 0, extra, 0)
        time.sleep(0.02)
        win32api.keybd_event(vk_code, 0, extra | win32con.KEYEVENTF_KEYUP, 0)
    except Exception:
        pass


class PPTController:
    """PPT 控制器：解析 PPTX + 键盘模拟翻页 + 文字匹配 + 幻灯片内高亮"""

    def __init__(self, ppt_path: str = ""):
        self.ppt_path = ppt_path
        # {页码: [(文字, left, top, width, height), ...]} 位置单位为 EMU
        self.page_texts_with_pos: dict[int, list[tuple[str, int, int, int, int]]] = {}
        self.page_texts: dict[int, list[str]] = {}  # 纯文字列表（向后兼容）
        self.total_pages = 0
        self.current_page = 0  # 0-based
        self._lock = threading.Lock()
        self._ppt_caption = os.path.basename(ppt_path) if ppt_path else ""
        # COM 相关
        self._ppt_app = None
        self._com_initialized = False
        # 屏幕尺寸缓存
        if win32api:
            self._screen_width = win32api.GetSystemMetrics(0)
            self._screen_height = win32api.GetSystemMetrics(1)
        else:
            self._screen_width = 1920
            self._screen_height = 1080

        # 演讲稿缓存 {页码(0-based): str}
        self.speech_notes: dict[int, str] = {}

        # 标记的重点页 {页码(0-based)}
        self.key_pages: set[int] = set()

        if ppt_path:
            self.load_pptx(ppt_path)
            self._print_start_instructions()

    def _init_com(self):
        """确保 COM 已初始化"""
        if not _HAS_COM:
            return False
        if not self._com_initialized:
            try:
                pythoncom.CoInitialize()
                self._com_initialized = True
            except Exception:
                return False
        return True

    def _get_ppt_app(self):
        """获取 PowerPoint COM 应用对象（带缓存）"""
        if not self._init_com():
            return None
        try:
            if self._ppt_app is None:
                self._ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
            return self._ppt_app
        except Exception:
            return None

    def load_pptx(self, ppt_path: str) -> bool:
        """加载并解析 PPTX 文件，建立文字+位置映射"""
        if Presentation is None:
            print("[PPT] python-pptx 未安装，无法解析 PPT 文字内容。")
            return False

        if not os.path.isfile(ppt_path):
            print(f"[PPT] 文件不存在: {ppt_path}")
            return False

        try:
            prs = Presentation(ppt_path)
            # 获取幻灯片尺寸（EMU）
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            self.page_texts_with_pos.clear()
            self.page_texts.clear()
            for slide_idx, slide in enumerate(prs.slides):
                texts_with_pos = []
                texts_only = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # 获取形状位置（EMU）
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts_with_pos.append((text, left, top, width, height))
                                texts_only.append(text)
                self.page_texts_with_pos[slide_idx] = texts_with_pos
                self.page_texts[slide_idx] = texts_only
            self.total_pages = len(prs.slides)
            self.current_page = 0
            print(f"[PPT] 已加载 {self.total_pages} 页，共 {sum(len(v) for v in self.page_texts.values())} 个文字块")
            return True
        except Exception as e:
            print(f"[PPT] 加载失败: {e}")
            return False

    def _print_start_instructions(self):
        """打印启动提示"""
        print("[PPT] ═══════════════════════════════════════════")
        print("[PPT]  请用 PowerPoint 或 WPS 手动打开此 PPT 文件")
        print(f"[PPT]  {self.ppt_path}")
        print("[PPT]  然后按 F5 开始全屏放映")
        print("[PPT]  手势/语音翻页会自动发送键盘指令")
        print("[PPT] ═══════════════════════════════════════════")

    def next_page(self):
        """下一页：键盘模拟（兼容任何放映软件）"""
        with self._lock:
            if self.current_page < self.total_pages - 1:
                self.current_page += 1
            page = self.current_page + 1
            total = self.total_pages
        _press_key(VK_NEXT)
        print(f"[PPT] → 第 {page}/{total} 页")
        self._try_com_sync_page()

    def prev_page(self):
        """上一页：键盘模拟（兼容任何放映软件）"""
        with self._lock:
            if self.current_page > 0:
                self.current_page -= 1
            page = self.current_page + 1
            total = self.total_pages
        _press_key(VK_PRIOR)
        print(f"[PPT] ← 第 {page}/{total} 页")
        self._try_com_sync_page()

    def _try_com_sync_page(self):
        """尝试通过 COM 同步实际页码（失败不影响使用）"""
        ppt_app = self._get_ppt_app()
        if ppt_app is None:
            return
        try:
            ss_window = ppt_app.ActivePresentation.SlideShowWindow
            if ss_window is not None:
                slide_index = ss_window.View.Slide.SlideIndex
                with self._lock:
                    self.current_page = slide_index - 1
        except Exception:
            pass

    def get_current_texts(self) -> list[str]:
        """获取当前页的所有文字块列表"""
        with self._lock:
            return self.page_texts.get(self.current_page, [])

    def get_current_page_number(self) -> int:
        """获取当前页码（1-based）"""
        with self._lock:
            return self.current_page + 1

    def get_total_pages(self) -> int:
        return self.total_pages

    def get_current_speech(self) -> str:
        """获取当前页的演讲稿"""
        with self._lock:
            return self.speech_notes.get(self.current_page, "")

    def set_speech_notes(self, notes: dict[int, str]):
        """设置演讲稿缓存"""
        with self._lock:
            self.speech_notes.clear()
            self.speech_notes.update(notes)
            page_count = len(notes)
            print(f"[PPT] 已加载 {page_count} 页演讲稿")

    def go_to_first_page(self):
        """跳到第一页"""
        with self._lock:
            self.current_page = 0
        _press_key(VK_HOME)
        print(f"[PPT] ⟹ 跳到首页 (第 1 页)")

    def go_to_last_page(self):
        """跳到最后一页"""
        with self._lock:
            self.current_page = self.total_pages - 1
        _press_key(VK_END)
        print(f"[PPT] ⟹ 跳到末页 (第 {self.total_pages} 页)")

    def toggle_black_screen(self):
        """黑屏/恢复（按 B 键）"""
        _press_key(VK_B)
        print("[PPT] ■ 黑屏切换")

    def quit_show(self):
        """结束放映（按 ESC 键）"""
        _press_key(VK_ESCAPE)
        print("[PPT] ■ 结束放映")

    def mark_current_page_as_key(self):
        """将当前页标记为重点页"""
        with self._lock:
            page = self.current_page
            if page not in self.key_pages:
                self.key_pages.add(page)
                print(f"[PPT] ★ 标记第 {page + 1} 页为重点页")
                return True
            return False

    def generate_key_pages_report(self) -> str:
        """生成重点页报告文本"""
        with self._lock:
            if not self.key_pages:
                return "未标记任何重点页。"

            lines = ["=" * 50]
            lines.append("📋 重点页报告")
            lines.append(f"PPT: {self._ppt_caption}")
            lines.append(f"共 {len(self.key_pages)} 个重点标记")
            lines.append("=" * 50)

            sorted_pages = sorted(self.key_pages)
            for idx, page in enumerate(sorted_pages, 1):
                page_num = page + 1  # 1-based
                texts = self.page_texts.get(page, [])
                # 取前 3 个文字块作为摘要
                summary = " | ".join(texts[:3]) if texts else "(无文字)"
                speech = self.speech_notes.get(page, "")
                speech_preview = (speech[:80] + "...") if len(speech) > 80 else speech if speech else "(无演讲稿)"

                lines.append(f"\n{idx}. 第 {page_num} 页")
                lines.append(f"   内容摘要: {summary}")
                lines.append(f"   演讲稿: {speech_preview}")

            lines.append("=" * 50)
            return "\n".join(lines)

    def match_text(self, spoken_text: str) -> list[str]:
        """
        将语音识别文本与当前页文字匹配
        返回匹配到的文字块列表
        """
        if not spoken_text or not spoken_text.strip():
            return []

        spoken_clean = re.sub(r'[^\w\u4e00-\u9fff]', '', spoken_text.lower())
        if not spoken_clean:
            return []

        matched = []
        current_texts = self.get_current_texts()
        for text in current_texts:
            text_clean = re.sub(r'[^\w\u4e00-\u9fff]', '', text.lower())
            if spoken_clean in text_clean or text_clean in spoken_clean:
                matched.append(text)

        return matched
