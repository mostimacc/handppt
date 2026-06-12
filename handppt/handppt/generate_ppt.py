# -*- coding: utf-8 -*-
"""生成课程汇报PPT（中文版）：基于MediaPipe手势识别与Vosk离线语音的智能PPT控制器

按照要求的7部分结构：
  1. 题目与应用场景
  2. 问题背景与任务定义
  3. 所选模型及原理简介
  4. 技术流程或实验设计
  5. 结果展示 / 案例分析 / 作品 Demo
  6. 问题反思与改进方向
  7. 总结学习收获
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x3C, 0x6E)
LIGHT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xE6, 0x8A, 0x2E)
BG_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)


def add_bg(slide, color=BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=BLACK, bold=False,
                align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, font_size=16, color=BLACK, bold=False,
             space_before=Pt(6), space_after=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def make_section_header(slide, number, title):
    """通用的节标题栏"""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.7),
                f'{number}  {title}', font_size=28, color=WHITE, bold=True)


# ====================================================================
# 第1页：封面
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)

add_textbox(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
            '基于MediaPipe手势识别与Vosk离线语音的',
            font_size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(1.2),
            '智能PPT控制系统',
            font_size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8),
            '——数字媒体技术中的多模态人机交互应用',
            font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
            '数字媒体技术 课程汇报',
            font_size=20, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# ====================================================================
# 第2页：目录
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
            '目  录', font_size=36, color=BLUE, bold=True)

sections = [
    ('01', '题目与应用场景'),
    ('02', '问题背景与任务定义'),
    ('03', '所选模型及原理简介'),
    ('04', '技术流程与实验设计'),
    ('05', '结果展示与案例分析'),
    ('06', '问题反思与改进方向'),
    ('07', '总结学习收获'),
]
for i, (num, title) in enumerate(sections):
    y = Inches(1.3) + Inches(i * 0.85)
    add_shape_bg(slide, Inches(1.2), y, Inches(0.8), Inches(0.6), BLUE)
    add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(0.8), Inches(0.5),
                num, font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.3), y + Inches(0.05), Inches(10), Inches(0.5),
                title, font_size=20, color=BLACK)

# ====================================================================
# 第3页：01 题目与应用场景
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '01', '题目与应用场景')

# 项目标题
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0), BG_LIGHT)
add_textbox(slide, Inches(1), Inches(1.4), Inches(11.3), Inches(0.8),
            '项目题目：基于MediaPipe手势识别与Vosk离线语音的智能PPT控制系统',
            font_size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# 场景说明
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
            '数字媒体应用场景', font_size=22, color=LIGHT_BLUE, bold=True)

scene_items = [
    '• 教室中：演讲者无需走回讲台按键盘，举手即可翻页',
    '• 会议室里：产品发布演示过程中，手势翻页 + 语音高亮关键信息',
    '• 大型演讲厅：自由走动演讲，零接触控制，提升沉浸感',
    '• 课堂教学：教师可同时用手势控制PPT和用激光笔指点内容',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(scene_items):
    p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=15, color=DARK_GRAY, space_before=Pt(6))

# 底部强调
add_shape_bg(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(1), Inches(5.7), Inches(11.3), Inches(0.6),
            '价值：手势控制 + 语音识别 = 多模态交互，让人更自然、更自由地操控数字内容',
            font_size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# 第4页：02 问题背景与任务定义
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '02', '问题背景与任务定义')

# 左栏：问题背景
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.5),
            '问题背景', font_size=22, color=LIGHT_BLUE, bold=True)
problem_items = [
    '传统PPT演示中，演讲者面临以下痛点：',
    '',
    '1. 需要频繁走回讲台按键翻页，打断演示节奏',
    '2. 翻页器功能单一，无法做标注或高亮',
    '3. 想展示某段文字时，需要鼠标定位滑动',
    '4. 演讲者活动范围受限，不能远离讲台',
    '5. 缺乏直观的多模态交互方式',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(problem_items):
    if item == '':
        add_para(tf, '', font_size=6, space_before=Pt(0))
    else:
        p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=14, color=DARK_GRAY, space_before=Pt(4))

# 右栏：任务定义
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
            '任务定义', font_size=22, color=LIGHT_BLUE, bold=True)
task_items = [
    '设计并实现一套"手势 + 语音"双模态PPT控制系统：',
    '',
    '▸ 手势控制模块',
    '   - 8种手势识别：滑动翻页、手指计数翻页、',
    '     OK切语音、大拇指首页、五指末页、握拳黑屏',
    '',
    '▸ 语音识别模块',
    '   - 离线中文语音→翻页/跳转/黑屏',
    '   - 语音匹配幻灯片文字 → 字幕显示',
    '',
    '▸ 系统要求',
    '   - 实时性：手势→翻页 < 150ms',
    '   - 兼容主流演示软件（PPT / WPS）',
    '   - 双模式动态切换（手势↔语音）',
]
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(5.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(task_items):
    if item == '':
        add_para(tf2, '', font_size=6, space_before=Pt(0))
    else:
        p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# ====================================================================
# 第5页：03 所选模型及原理简介
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '03', '所选模型及原理简介')

# 模型1：MediaPipe Hand Landmarker
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.5),
            '模型一：MediaPipe Hand Landmarker', font_size=20, color=LIGHT_BLUE, bold=True)
mp_items = [
    '类型：轻量级CNN + 热力图回归',
    '',
    '▸ 输入：640×480 RGB图像帧（视频流）',
    '▸ 输出：21个3D手部关键点 (x, y, z) + 置信度',
    '',
    '▸ 为什么选择该模型：',
    '   - Google MediaPipe团队专为移动端/实时场景优化',
    '   - GPU推理仅需 < 30ms，CPU < 80ms',
    '   - 提供手部 21 点拓扑结构，支持灵活的手势规则判断',
    '',
    '▸ 核心原理：先检测手掌区域（轻量BlazePalm），',
    '   再回归21个关键点在手掌坐标系中的位置',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(mp_items):
    if item == '':
        add_para(tf, '', font_size=6, space_before=Pt(0))
    else:
        p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# 模型2：Vosk
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
            '模型二：Vosk离线语音识别', font_size=20, color=LIGHT_BLUE, bold=True)
vk_items = [
    '类型：基于Kaldi的DNN-HMM混合模型',
    '       核心：CNN（时延神经网络）+ LSTM + HMM',
    '',
    '▸ 输入：16kHz单声道PCM音频流',
    '▸ 输出：识别文本（final结果 + partial中间结果）',
    '',
    '▸ 为什么选择该模型：',
    '   - 完全离线运行，保护演讲隐私',
    '   - 中文识别准确率 > 85%（安静环境）',
    '   - 支持流式推理（回调返回中间结果）',
    '',
    '▸ 核心原理：',
    '   CNN提取声学特征 → LSTM建模时序依赖 → '
    'HMM状态解码 → 语言模型输出文字',
]
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(vk_items):
    if item == '':
        add_para(tf2, '', font_size=6, space_before=Pt(0))
    else:
        p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# 底部总结
add_shape_bg(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8), BG_LIGHT)
add_textbox(slide, Inches(1), Inches(6.45), Inches(11.3), Inches(0.7),
            '选型逻辑：MediaPipe提供实时手部关键点 → 手势规则转化；Vosk提供离线语音 → 文本匹配。'
            '二者形成互补的双模态控制方案，覆盖"举手翻页"和"动口控制"两种场景。',
            font_size=14, color=BLUE)

# ====================================================================
# 第6页：04 技术流程与实验设计
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '04', '技术流程与实验设计')

# 数据来源 & 预处理
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.5),
            '数据来源与预处理', font_size=20, color=LIGHT_BLUE, bold=True)
data_items = [
    '视频数据：USB摄像头实时采集 640×480 @ 30fps',
    '  → BGR转RGB色彩空间 → 送入MediaPipe管道',
    '',
    '音频数据：麦克风采集 16kHz 单声道PCM',
    '  → 静音检测（SilenceTimeout=2s）→ Vosk流式引擎',
    '',
    'PPT数据：python-pptx解析.pptx文件',
    '  → 提取每页幻灯片文字内容，建立 页码→[文字]映射',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(data_items):
    if item == '':
        add_para(tf, '', font_size=6, space_before=Pt(0))
    else:
        p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# 系统流程
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
            '系统处理流程', font_size=20, color=LIGHT_BLUE, bold=True)
pipeline_items = [
    '手势控制模式（默认进入）：',
    '  摄像头 → MediaPipe → 几何规则（8种手势）',
    '    → 稳定性校验（连续8帧稳定）→ 模拟键盘翻页',
    '',
    '语音识别模式（OK手势进入）：',
    '  麦克风 → Vosk流式识别 → 关键词匹配',
    '    → 翻页/跳转/黑屏/文字匹配 → 字幕显示',
    '',
    '模式切换机制：',
    '  OK手势 → 启动语音监听 → 进入语音模式',
    '  说"ok"/"好的" → 停止语音 → 回到手势模式',
]
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(pipeline_items):
    if item == '':
        add_para(tf2, '', font_size=6, space_before=Pt(0))
    else:
        p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# 底部：实验设计
add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.8), BG_LIGHT)
add_textbox(slide, Inches(1), Inches(5.4), Inches(11.3), Inches(0.4),
            '实验设计要点', font_size=16, color=BLUE, bold=True)
exp_items = [
    '• 测试环境：Windows 11 + PowerPoint/WPS，摄像头帧率 30fps，不同光照条件',
    '• 评价指标：手势识别准确率、语音识别准确率、系统端到端延迟',
    '• 测试用例：8种手势各测试20次，10句语音命令各测试10次，记录成功/失败',
]
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(1.2))
tf3 = txBox3.text_frame
tf3.word_wrap = True
for i, item in enumerate(exp_items):
    p = tf3.paragraphs[0] if i == 0 else add_para(tf3, item, font_size=13, color=DARK_GRAY, space_before=Pt(4))

# ====================================================================
# 第7页：05 结果展示与案例分析
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '05', '结果展示与案例分析')

# 左栏：功能结果
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.5),
            '功能验证结果', font_size=20, color=LIGHT_BLUE, bold=True)
result_items = [
    '手势识别（8种手势）：',
    '   ✓ 手掌右滑 → 下一页（准确率 > 90%）',
    '   ✓ 手掌左滑 → 上一页（准确率 > 90%）',
    '   ✓ OK手势   → 切语音（准确率 > 95%）',
    '   ✓ 伸1指/2指 → 翻页（补充手势）',
    '   ✓ 大拇指/五指/握拳 → 首页/末页/黑屏',
    '',
    '语音识别：',
    '   ✓ 翻页命令（下一页/上一页）',
    '   ✓ 跳转命令（首页/末页）',
    '   ✓ 语音匹配PPT文字 → 字幕显示PPT原文',
    '',
    '模式切换：OK → 语音模式 / 说"ok" → 手势模式',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(result_items):
    if item == '':
        add_para(tf, '', font_size=6, space_before=Pt(0))
    else:
        p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# 右栏：评价指标 + Demo
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
            '评价指标与Demo展示', font_size=20, color=LIGHT_BLUE, bold=True)
eval_items = [
    '实时性指标：',
    '   • 手势→翻页：< 150ms（含8帧稳定性校验）',
    '   • 语音→匹配：< 500ms（流式中间结果）',
    '   • 单帧处理：< 80ms（CPU模式）',
    '',
    '准确率指标：',
    '   • OK手势识别准确率：> 95%',
    '   • 左右滑动手势：> 90%',
    '   • 中文语音命令：> 85%（安静环境）',
    '',

    '兼容性：',
    '   • Microsoft PowerPoint ✓',
    '   • WPS Office ✓',
    '',
    '现场Demo展示：实时摄像头画面 → 手势翻页',
    '→ 用语音命令翻回上一页',
    '→ 说"ok"切回手势 → 翻到末页',
]
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(5.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(eval_items):
    if item == '':
        add_para(tf2, '', font_size=6, space_before=Pt(0))
    else:
        p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# ====================================================================
# 第8页：06 问题反思与改进方向
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_section_header(slide, '06', '问题反思与改进方向')

# 左栏：当前不足
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.5),
            '当前方案的不足', font_size=20, color=LIGHT_BLUE, bold=True)
short_items = [
    '1. 仅支持8种预定义手势，无法识别动态连续动作',
    '   （如挥手、画圈示意等复杂姿态）',
    '',
    '2. Vosk在嘈杂环境（如会议室多人讨论）',
    '   下准确率明显下降',
    '',
    '3. 匹配PPT文字依赖python-pptx预解析，',
    '   不支持图片内嵌文字或SmartArt中的文字',
    '',
    '4. 摄像头固定，演讲者活动范围受限',
    '',
    '5. 缺少多人场景下的指定主讲人支持',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(short_items):
    if item == '':
        add_para(tf, '', font_size=6, space_before=Pt(0))
    else:
        p = tf.paragraphs[0] if i == 0 else add_para(tf, item, font_size=14, color=DARK_GRAY, space_before=Pt(4))

# 右栏：数据偏差/风险 + 优化方向
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
            '数据偏差、风险与优化', font_size=20, color=LIGHT_BLUE, bold=True)
risk_items = [
    '▸ 数据偏差：',
    '   - 测试多为右手操作，缺少左撇子用户数据',
    '   - 手势阈值基于特定光照条件调参',
    '',
    '▸ 隐私与安全：',
    '   - 摄像头画面本地处理（未加密传输）',
    '   - 建议增加用户知情同意提示',
    '',
    '▸ 模型可解释性：',
    '   - MediaPipe为端到端黑盒模型',
    '   - 手势分类依赖人工几何规则而非学习',
    '',
    '▸ 后续优化方向：',
    '   - 引入Gesture Transformer识别连续手势',
    '   - 集成Whisper增强噪声鲁棒性',
    '   - 增加手势激光笔/涂鸦功能',
    '   - 支持多手追踪 + 主讲人区分',
]
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(5.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(risk_items):
    if item == '':
        add_para(tf2, '', font_size=6, space_before=Pt(0))
    else:
        p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=13, color=DARK_GRAY, space_before=Pt(3))

# ====================================================================
# 第9页：07 总结学习收获
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
            '07  总结学习收获', font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 项目成果总结
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.0), RGBColor(0x1E, 0x50, 0x8A))
add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.4),
            '项目成果', font_size=20, color=ACCENT, bold=True)
summary_items = [
    '构建了一套完整的手势 + 语音双模态PPT控制系统',
    '实现了8种手势的实时识别与稳定触发',
    '集成了离线中文语音识别，支持多关键词命令与PPT文字匹配',
    '实现了手势↔语音动态无缝切换，演讲者自由选择控制方式',
    '系统兼容Microsoft PowerPoint和WPS两大主流演示平台',
]
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.3))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(summary_items):
    p = tf.paragraphs[0] if i == 0 else add_para(tf, f'• {item}', font_size=14, color=RGBColor(0xDD, 0xDD, 0xDD), space_before=Pt(4))

# 学习收获
add_shape_bg(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(2.5), RGBColor(0x1E, 0x50, 0x8A))
add_textbox(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.4),
            '学习收获', font_size=20, color=ACCENT, bold=True)
learn_items = [
    '• 深入理解了MediaPipe轻量级CNN手部检测的原理与局限',
    '• 掌握了Vosk离线语音识别模型的部署与流式调用方法',
    '• 实践了多模态人机交互系统的架构设计（双线程+队列+模式管理）',
    '• 体会到数字媒体技术中"计算机视觉+语音处理+人机交互"的融合价值',
    '• 认识到现实部署中的挑战：光照变化、噪声干扰、用户习惯差异',
]
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(learn_items):
    p = tf2.paragraphs[0] if i == 0 else add_para(tf2, item, font_size=14, color=RGBColor(0xDD, 0xDD, 0xDD), space_before=Pt(4))

# 致谢
add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
            '感谢聆听！欢迎提问。', font_size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
            '数字媒体技术 课程汇报  |  基于MediaPipe + Vosk 的智能PPT控制系统',
            font_size=14, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# 保存
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, '智能PPT控制器_课程汇报.pptx')
prs.save(output_path)
print(f'PPT已保存至: {output_path}')
print(f'共 {len(prs.slides)} 页幻灯片')